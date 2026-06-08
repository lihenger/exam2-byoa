"""PPTX 文本提取模块"""

import os
from pptx import Presentation


def read_pptx(path: str) -> dict:
    """读取 PPTX 文件，返回按幻灯片组织的结构化内容。"""
    if not os.path.exists(path):
        raise FileNotFoundError(f"文件不存在: {path}")

    prs = Presentation(path)
    slides = []
    total_words = 0

    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    t = paragraph.text.strip()
                    if t:
                        texts.append(t)
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    row_texts = [cell.text.strip() for cell in row.cells]
                    texts.append(" | ".join(row_texts))

        slide_text = "\n".join(texts)
        wc = len(slide_text.split())
        total_words += wc
        slides.append({
            "slide_num": i + 1,
            "text": slide_text,
            "element_count": len(texts),
            "word_count": wc,
        })

    return {
        "file_path": os.path.abspath(path),
        "file_name": os.path.basename(path),
        "file_type": "pptx",
        "total_slides": len(slides),
        "total_words": total_words,
        "slides": slides,
    }
